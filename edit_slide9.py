import os
import copy
from lxml import etree
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor

home = os.path.expanduser('~')
src = os.path.join(home, 'Downloads', 'ZizoCircle Pitch v9.pptx')
# Output goes to the uploads path referenced in the task
out = os.path.join(home, 'Downloads', 'ZizoCircle_Pitch_v10.pptx')

prs = Presentation(src)
slide = prs.slides[8]  # index 8 = slide 9

# ------------------------------------------------------------------ #
# Helper: set a single-run paragraph's text preserving run formatting #
# ------------------------------------------------------------------ #
def set_single_run_text(shape, new_text):
    """Replace text in a shape that has exactly 1 paragraph with 1 run."""
    tf = shape.text_frame
    para = tf.paragraphs[0]
    run = para.runs[0]
    old = run.text
    run.text = new_text
    print(f"  [CHANGED] '{old}' -> '{new_text}'")

# ------------------------------------------------------------------ #
# Build a lookup of shapes by name for easy access                    #
# ------------------------------------------------------------------ #
shapes_by_name = {s.name: s for s in slide.shapes}

# ------------------------------------------------------------------ #
# MILESTONE DATA                                                       #
# 8 milestones: (label, description, colour_hex)                      #
# Rows 1-3 keep green (27C987), rows 4-8 use purple (C94FC9)          #
# Row 3 and 4 share "April 2026" — both get green for label           #
# ------------------------------------------------------------------ #
milestones = [
    ("Q3 2025",           "Concept, Business Plan & UI/UX Design Finalized",          "27C987", "FFFFFF"),
    ("Q4 2025",           "App Development Started",                                   "27C987", "FFFFFF"),
    ("April 2026",        "MVP Built & Live — Android APK available for download",     "27C987", "FFFFFF"),
    ("April 2026",        "First real users registered — Firebase backend live",       "27C987", "FFFFFF"),
    ("May 2026",          "Snoonu Incubator Application — seeking QAR 1,255,000",      "C94FC9", "E8E0F7"),
    ("45 Days\npost-funding", "Full public launch iOS & Android + marketing begins",   "C94FC9", "E8E0F7"),
    ("Month 2-4",         "100+ Partner Venues Onboarded · Influencer Campaign Active","C94FC9", "E8E0F7"),
    ("Month 6-12",        "50K+ Users · QAR 300K+ Monthly Revenue · UAE Expansion Prep","C94FC9","E8E0F7"),
]

# ------------------------------------------------------------------ #
# Existing milestone shape groups (label_shape_name, desc_shape_name, #
# dot_shape_name)                                                      #
# ------------------------------------------------------------------ #
existing_rows = [
    ("Text 3",  "Text 4",  "Shape 2"),
    ("Text 6",  "Text 7",  "Shape 5"),
    ("Text 9",  "Text 10", "Shape 8"),
    ("Text 12", "Text 13", "Shape 11"),
    ("Text 15", "Text 16", "Shape 14"),
    ("Text 18", "Text 19", "Shape 17"),
]

# ------------------------------------------------------------------ #
# Vertical layout: we need 8 rows.                                    #
# Current rows span from top=960120 to bottom of row6 ~4526280        #
# Available vertical space (excluding status bar) ~ 3566160 EMU       #
# We'll space rows evenly starting at same top as row1                #
# ------------------------------------------------------------------ #
TOP_FIRST  = 960120     # top of first milestone label (EMU)
# Status bar background starts at 4617720; leave ~120K gap
AVAIL_H    = 4617720 - 120000 - TOP_FIRST  # ≈3537600 EMU
ROW_COUNT  = 8
ROW_STEP   = AVAIL_H // (ROW_COUNT - 1)    # spacing between row tops

print(f"\nRow step: {ROW_STEP} EMU = {ROW_STEP/914400*2.54:.2f} cm")

# Timeline bar (Shape 1): stretch it to cover all 8 rows
timeline_bar = shapes_by_name["Shape 1"]
new_bar_top = TOP_FIRST + (ROW_STEP // 2)           # centre on first dot
last_row_top = TOP_FIRST + (ROW_COUNT - 1) * ROW_STEP
new_bar_bottom = last_row_top + (ROW_STEP // 2)     # centre on last dot
old_bar_top    = timeline_bar.top
old_bar_height = timeline_bar.height
timeline_bar.top    = new_bar_top
timeline_bar.height = new_bar_bottom - new_bar_top
print(f"\n[CHANGED] Timeline bar top: {old_bar_top} -> {new_bar_top}")
print(f"[CHANGED] Timeline bar height: {old_bar_height} -> {timeline_bar.height}")

# ------------------------------------------------------------------ #
# Helper: deep-copy a shape's XML element and add it to the slide     #
# ------------------------------------------------------------------ #
def clone_shape_to_slide(source_shape, slide):
    sp = copy.deepcopy(source_shape._element)
    slide.shapes._spTree.append(sp)
    # Return the newly added shape (last one)
    return slide.shapes[-1]

# ------------------------------------------------------------------ #
# Update existing 6 rows and create 2 new ones                        #
# ------------------------------------------------------------------ #

# Template shapes for cloning (row 1)
tmpl_label = shapes_by_name["Text 3"]
tmpl_desc  = shapes_by_name["Text 4"]
tmpl_dot   = shapes_by_name["Shape 2"]

new_shapes = []  # will hold (label_shape, desc_shape, dot_shape) for all 8 rows

for row_idx, (label_text, desc_text, label_color_hex, desc_color_hex) in enumerate(milestones):
    row_top = TOP_FIRST + row_idx * ROW_STEP
    dot_top = row_top + (365760 // 2) - (310896 // 2)  # vertically centre dot on label

    if row_idx < 6:
        # Use existing shapes
        lname, dname, dotname = existing_rows[row_idx]
        lshape  = shapes_by_name[lname]
        dshape  = shapes_by_name[dname]
        dotshape = shapes_by_name[dotname]

        # Reposition
        lshape.top   = row_top
        dshape.top   = row_top
        dotshape.top = dot_top

        # Update label text + colour
        lpara = lshape.text_frame.paragraphs[0]
        lrun  = lpara.runs[0]
        old_l = lrun.text
        lrun.text = label_text
        lrun.font.color.rgb = RGBColor.from_string(label_color_hex)
        print(f"\nRow {row_idx+1} label  [{lname}]: '{old_l}' -> '{label_text}' (color {label_color_hex})")

        # Update desc text + colour
        dpara = dshape.text_frame.paragraphs[0]
        drun  = dpara.runs[0]
        old_d = drun.text
        drun.text = desc_text
        drun.font.color.rgb = RGBColor.from_string(desc_color_hex)
        print(f"Row {row_idx+1} desc   [{dname}]: '{old_d}' -> '{desc_text}' (color {desc_color_hex})")

    else:
        # Clone from existing row 1 templates (index 0 = row 1)
        print(f"\nRow {row_idx+1}: CREATING NEW shapes by cloning templates")

        new_label = clone_shape_to_slide(tmpl_label, slide)
        new_desc  = clone_shape_to_slide(tmpl_desc,  slide)
        new_dot   = clone_shape_to_slide(tmpl_dot,   slide)

        # Reposition
        new_label.left = tmpl_label.left
        new_label.top  = row_top
        new_label.width  = tmpl_label.width
        new_label.height = tmpl_label.height

        new_desc.left  = tmpl_desc.left
        new_desc.top   = row_top
        new_desc.width   = tmpl_desc.width
        new_desc.height  = tmpl_desc.height

        new_dot.left  = tmpl_dot.left
        new_dot.top   = dot_top
        new_dot.width   = tmpl_dot.width
        new_dot.height  = tmpl_dot.height

        # Set label text + colour
        lrun = new_label.text_frame.paragraphs[0].runs[0]
        lrun.text = label_text
        lrun.font.color.rgb = RGBColor.from_string(label_color_hex)
        print(f"Row {row_idx+1} label  [NEW]: '{label_text}' (color {label_color_hex})")

        # Set desc text + colour
        drun = new_desc.text_frame.paragraphs[0].runs[0]
        drun.text = desc_text
        drun.font.color.rgb = RGBColor.from_string(desc_color_hex)
        print(f"Row {row_idx+1} desc   [NEW]: '{desc_text}' (color {desc_color_hex})")

# ------------------------------------------------------------------ #
# Update status bar (Shape 21)                                        #
# ------------------------------------------------------------------ #
status_shape = shapes_by_name["Text 21"]
status_para  = status_shape.text_frame.paragraphs[0]
status_run   = status_para.runs[0]
old_status   = status_run.text
new_status   = "Stage: MVP LIVE · Android APK available · First users onboarded · Firebase backend active · Seeking: QAR 1,255,000"
status_run.text = new_status
print(f"\n[CHANGED] Status bar: '{old_status}' -> '{new_status}'")

# ------------------------------------------------------------------ #
# Save                                                                #
# ------------------------------------------------------------------ #
prs.save(out)
print(f"\n✓ Saved to: {out}")
