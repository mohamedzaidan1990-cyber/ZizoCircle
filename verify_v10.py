import os
from pptx import Presentation

home = os.path.expanduser('~')
out = os.path.join(home, 'Downloads', 'ZizoCircle_Pitch_v10.pptx')

prs = Presentation(out)
slide = prs.slides[8]

print("=== SLIDE 9 v10 VERIFICATION ===")
print(f"Total shapes: {len(slide.shapes)}\n")

for i, shape in enumerate(slide.shapes):
    if shape.has_text_frame and shape.text_frame.text.strip():
        print(f"Shape {i} [{shape.name}]: '{shape.text_frame.text}' | top={shape.top/914400*2.54:.2f}cm")

print("\n=== OTHER SLIDES CHECK (slide 9 = index 8 only changed) ===")
for si, s in enumerate(prs.slides):
    if si == 8:
        continue
    for shape in s.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text
            if 'MVP LIVE' in txt or 'Firebase backend active' in txt:
                print(f"  WARNING: slide {si+1} shape '{shape.name}' contains modified text!")
print("Other slides check complete.")
