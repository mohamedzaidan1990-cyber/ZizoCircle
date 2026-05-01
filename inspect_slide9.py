from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

import os
home = os.path.expanduser('~')
src = os.path.join(home, 'Downloads', 'ZizoCircle Pitch v9.pptx')
prs = Presentation(src)

slide = prs.slides[8]  # index 8 = slide 9

print("=== SLIDE 9 INSPECTION ===")
print(f"Total shapes: {len(slide.shapes)}\n")

for i, shape in enumerate(slide.shapes):
    print(f"--- Shape {i} ---")
    print(f"  Name: {shape.name}")
    print(f"  Shape type: {shape.shape_type}")
    print(f"  Left: {shape.left}, Top: {shape.top}, Width: {shape.width}, Height: {shape.height}")
    print(f"  Left(cm): {shape.left/914400*2.54:.2f}, Top(cm): {shape.top/914400*2.54:.2f}")
    print(f"  Width(cm): {shape.width/914400*2.54:.2f}, Height(cm): {shape.height/914400*2.54:.2f}")

    if shape.has_text_frame:
        tf = shape.text_frame
        print(f"  Text frame paragraphs: {len(tf.paragraphs)}")
        for pi, para in enumerate(tf.paragraphs):
            print(f"  Para {pi}: '{para.text}'")
            for ri, run in enumerate(para.runs):
                font = run.font
                color_str = "None"
                if font.color and font.color.type:
                    try:
                        color_str = str(font.color.rgb)
                    except Exception:
                        color_str = f"type={font.color.type}"
                print(f"    Run {ri}: '{run.text}' | size={font.size} | bold={font.bold} | italic={font.italic} | color={color_str} | name={font.name}")
    elif hasattr(shape, 'text'):
        print(f"  Text: {shape.text[:200]}")
    print()
